import re
import pptx.exc
from pptx.enum.shapes import *
from pptx import *


def fractional_in_text_checker(text: str) -> list[str]:
    """
    Функция для поиска дробных чисел в тексте.
    :param text: Входной текст.
    :return: Список, содержащий все дробные числа из текста.
    """
    special_fraction_symbols = "¼½¾"
    pattern = re.compile(r'(?<![\d.,-])-?(?:\d+[,.]\d+|\d+/\d+|[¼½¾])(?!\.\d)(?!,\d)(?!\d)')
    numbers = pattern.findall(text.replace(",", "."))
    return [num for num in numbers if not num.endswith((".0", ",0")) or num in special_fraction_symbols]


def chart_handler(chart_shape: MSO_SHAPE_TYPE.CHART) -> str:
    """
    Функция для извлечения текста из диаграммы.
    :param chart_shape: Объект, представляющий диаграмму.
    :return: Строка со значениями всех полей диаграммы
    """
    text = ''
    for series in chart_shape.chart.series:
        for value in series.values:
            text += f" {value}"
    return text


def table_handler(table_shape: MSO_SHAPE_TYPE.TABLE) -> str:
    """
    Функция для извлечения текста из таблицы.
    :param table_shape: Объект, представляющий таблицу.
    :return: Строка со значениями всех ячеек таблицы
    """
    text = ''
    for row in table_shape.table.rows:
        for cell in row.cells:
            text += f" {cell.text}"
    return text


def slide_processing(slide: Presentation().slides, slide_number: int) -> list[int | list[str]]:
    """
    Обрабатывает один слайд, извлекая из него дробные числа.
    :param slide: Объект слайда.
    :param slide_number: Номер слайда.
    :return: Список, содержащий номер слайда и список дробных чисел.
    """
    all_num = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                numbers = fractional_in_text_checker(shape.text)
                if numbers:
                    all_num += numbers
            else:
                match shape.shape_type:
                    case MSO_SHAPE_TYPE.TABLE:
                        numbers = fractional_in_text_checker(table_handler(shape))
                        if numbers:
                            all_num += numbers
                    case MSO_SHAPE_TYPE.CHART:
                        numbers = fractional_in_text_checker(chart_handler(shape))
                        if numbers:
                            all_num += numbers
                    case _:
                        continue
        except Exception as err:
            print(f"Ошибка обработчика слайдов: {err}")
    fractional_numbers = [slide_number, all_num]
    return fractional_numbers


def search_for_fractional_numbers_in_pptx(path_to_pptx: str) -> dict[int: str]:
    """
    Функция для поиска дробных чисел по всей презентации.
    :param path_to_pptx: Путь к файлу .pptx.
    :return: Словарь, где ключ - номер слайда, значение - список дробных чисел, или {}, если ничего не найдено.
    """
    try:
        if not (path_to_pptx.endswith(".pptx")):
            path_to_pptx += ".pptx"
        try:
            all_fractional_numbers_from_pptx = {}
            prs = Presentation(path_to_pptx)
            slide_number = 0
            for slide in prs.slides:
                slide_number += 1
                from_slide_fractnum = slide_processing(slide, slide_number)
                if len(from_slide_fractnum[1]) > 0:
                    all_fractional_numbers_from_pptx[from_slide_fractnum[0]] = from_slide_fractnum[1]
            if not all_fractional_numbers_from_pptx:
                print("Дробных чисел в презентации не найдено")
                return {}
            return all_fractional_numbers_from_pptx
        except pptx.exc.PackageNotFoundError:
            print(
                "По указанному пути презентации .pptx не найдено.\nПроверьте правильность указания ссылки и тип файла")
            return {}

    except Exception as e:
        print(f"Произошла непредвиденная ошибка: {e}")
        return {}


if __name__ == '__main__':
    path = input("Введите путь до .pptx презентации: ")
    result = search_for_fractional_numbers_in_pptx(path)
    for slide_number, numbers in result.items():
        print(f"\nСлайд номер {slide_number}:")
        for number in numbers:
            print(number)
